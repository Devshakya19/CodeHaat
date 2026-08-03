#!/usr/bin/env python3
"""
Convert markdown presentation to PowerPoint (.pptx) file
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re
import sys


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    # Format title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = 'Calibri'

    # Format subtitle
    subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle_placeholder.text_frame.paragraphs[0].font.italic = True
    subtitle_placeholder.text_frame.paragraphs[0].font.name = 'Calibri'

    return slide


def add_content_slide(prs, title, content_lines):
    """Add a content slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(32)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = 'Calibri'

    # Process content lines
    text_frame = body_placeholder.text_frame
    text_frame.clear()  # Clear existing content

    for i, line in enumerate(content_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        # Check if it's a bullet point
        if line.strip().startswith('-'):
            p.text = line.strip()[1:].strip()
            p.level = 0  # Bullet level
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
        elif line.strip().startswith('*'):
            p.text = line.strip()[1:].strip()
            p.level = 0  # Bullet level
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
        elif re.match(r'^\d+\.', line.strip()):  # Numbered list
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
        else:
            p.text = line.strip()
            p.level = 0
            p.font.size = Pt(18)
            p.font.name = 'Calibri'

        # Add space after paragraph
        p.space_after = Pt(6)

    return slide


def add_section_header_slide(prs, title):
    """Add a section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title

    title_placeholder.text = title
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
    title_placeholder.text_frame.paragraphs[0].font.bold = True
    title_placeholder.text_frame.paragraphs[0].font.name = 'Calibri'
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)  # Blue

    return slide


def parse_markdown(file_path):
    """Parse markdown file into sections"""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by horizontal rules
    sections = re.split(r'\n---\n', content)

    parsed_sections = []
    for section in sections:
        if not section.strip():
            continue

        lines = section.strip().split('\n')
        title = ""
        content_lines = []

        for line in lines:
            # Check for main section headers (like "# 1. Executive Summary")
            if re.match(r'^#\s+\d+\. ', line):
                if title:  # If we already have a title, save previous section
                    parsed_sections.append((title, content_lines))
                    title = line[2:].strip()  # Remove "# "
                    content_lines = []
                else:
                    title = line[2:].strip()  # Remove "# "
            # Check for subsection headers (like "## The Opportunity")
            elif line.startswith('## '):
                if title:  # If we already have a title, save previous section
                    parsed_sections.append((title, content_lines))
                    title = line[3:].strip()  # Remove "## "
                    content_lines = []
                else:
                    title = line[3:].strip()  # Remove "## "
            # Check for regular content lines
            elif line.strip() and not line.startswith('>') and not line.strip() == '':
                # Skip table of contents lines that are just numbers
                if re.match(r'^\d+\.\s*[A-Z]', line.strip()):
                    continue
                content_lines.append(line)

        # Don't forget the last section
        if title or content_lines:
            if not title:
                title = "Content"
            parsed_sections.append((title, content_lines))

    return parsed_sections


def create_presentation_from_markdown(md_path, pptx_path):
    """Create a PowerPoint presentation from markdown file"""
    # Parse markdown
    sections = parse_markdown(md_path)

    # Create presentation
    prs = Presentation()
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Add title slide
    title_slide = add_title_slide(prs, "CodeHaat — Investor & Partner Presentation", 
                                 "India's #1 Digital Code Marketplace — Where Code Meets Commerce")

    # Add table of contents slide
    toc_content = [
        "1. Executive Summary",
        "2. Problem Statement", 
        "3. Our Solution",
        "4. Technology & Innovation",
        "5. Market Opportunity",
        "6. Business Model & Traction",
        "7. Go-to-Market Strategy",
        "8. Competitive Analysis",
        "9. Financial Projections",
        "10. Team & Advisors",
        "11. Roadmap & Milestones",
        "12. Investment Ask",
        "13. Contact Information"
    ]
    add_content_slide(prs, "Table of Contents", toc_content)

    # Add sections
    for title, content_lines in sections:
        # Skip if it's just the title or table of contents
        if title.lower() in ["table of contents", "codehaat — investor & partner presentation"]:
            continue
            
        # Add section header
        add_section_header_slide(prs, title)
        
        # Add content slide if there's content
        if content_lines:
            add_content_slide(prs, title, content_lines[:10])  # Limit to 10 bullet points per slide

    # Save presentation
    prs.save(pptx_path)
    print(f"Presentation saved to {pptx_path}")


if __name__ == "__main__":
    md_file = "/home/ghost/Projects/CodeHaat/docs/08-PRESENTATION.md"
    pptx_file = "/home/ghost/Projects/CodeHaat/docs/08-PRESENTATION.pptx"
    
    create_presentation_from_markdown(md_file, pptx_file)